import os
import uuid
import shutil
import zipfile # Added for zipping images from PDF
from flask import Flask, render_template, request, redirect, url_for, send_from_directory, flash, jsonify
from werkzeug.utils import secure_filename
from PIL import Image

# Core Conversion/Merging Libraries
from PyPDF2 import PdfReader, PdfWriter, PdfMerger # PyPDF2 3.0+
from docx import Document as DocxDocument
from pptx import Presentation
from docx2pdf import convert as docx2pdf_convert
from pdf2image import convert_from_path as pdf_to_images_convert # Needs poppler
from pdf2docx import Converter as Pdf2DocxConverter
import img2pdf

from reportlab.pdfgen import canvas as reportlab_canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader as ReportlabImageReader


app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "a_very_strong_default_secret_key_that_should_be_changed_for_prod")

# --- Configuration ---
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
CONVERTED_FOLDER = os.path.join(BASE_DIR, 'converted')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

ALLOWED_UPLOAD_EXTENSIONS = {'pdf', 'docx', 'pptx', 'jpg', 'jpeg', 'png'}
IMAGE_EXTENSIONS = {'jpg', 'jpeg', 'png'}

def get_file_ext(filename): # Same helper
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

def cleanup_files(*filepaths): # Same helper
    for filepath in filepaths:
        # ... (same as previous)
        if filepath and os.path.exists(filepath):
            try:
                if os.path.isdir(filepath):
                    shutil.rmtree(filepath)
                else:
                    os.remove(filepath)
                app.logger.info(f"Cleaned up: {filepath}")
            except Exception as e:
                app.logger.error(f"Error cleaning up file/folder {filepath}: {e}")


# --- Conversion Functions (These should be the robust versions from previous enhanced prompt) ---
def convert_pdf_to_docx_robust(input_path, output_path): # Same
    try:
        cv = Pdf2DocxConverter(input_path)
        cv.convert(output_path)
        cv.close()
        return True, "PDF converted to DOCX successfully."
    except Exception as e:
        app.logger.error(f"pdf2docx conversion error: {e}", exc_info=True)
        return False, f"PDF to DOCX conversion failed. Error: {str(e)}. Complex PDFs may not convert well."

def convert_docx_to_pdf_robust(input_path, output_path): # Same
    try:
        docx2pdf_convert(input_path, output_path)
        return True, "DOCX converted to PDF successfully."
    except Exception as e:
        app.logger.error(f"docx2pdf conversion error: {e}", exc_info=True)
        return False, f"DOCX to PDF conversion failed. Ensure LibreOffice or MS Office is available. Error: {str(e)}"

def convert_images_to_pdf_robust(image_paths, output_path): # Same
    try:
        valid_image_paths = [p for p in image_paths if get_file_ext(os.path.basename(p)) in IMAGE_EXTENSIONS]
        if not valid_image_paths:
            return False, "No valid images provided for PDF conversion."
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(valid_image_paths))
        return True, "Images converted to PDF successfully."
    except Exception as e:
        app.logger.error(f"Images to PDF conversion error: {e}", exc_info=True)
        return False, f"Failed to convert images to PDF. Error: {str(e)}"

def convert_pdf_to_images_robust(input_path, output_folder_stem_for_naming): # Same, returns path to zip or single image
    # output_folder_stem_for_naming is like "converted/uniqueid_output"
    output_image_paths = []
    temp_image_dir_name = f"{os.path.basename(output_folder_stem_for_naming)}_img_pages_{uuid.uuid4().hex[:4]}"
    temp_image_dir_path = os.path.join(app.config['CONVERTED_FOLDER'], temp_image_dir_name) # Create temp dir inside converted

    try:
        os.makedirs(temp_image_dir_path, exist_ok=True)
        # pdf2image saves files directly, output_file is a prefix for those files
        # it returns a list of PIL Image objects, not paths when output_folder is used.
        # Let's save them manually to have control over paths.
        images_pil = pdf_to_images_convert(input_path, fmt='png') # poppler_path can be specified if not in PATH
        
        if not images_pil:
             return False, "No images could be extracted from the PDF.", None

        for i, img_pil in enumerate(images_pil):
            img_filename = f"{os.path.basename(output_folder_stem_for_naming)}_page_{i+1}.png"
            img_path = os.path.join(temp_image_dir_path, img_filename)
            img_pil.save(img_path, "PNG")
            output_image_paths.append(img_path)
        
        if not output_image_paths: # Should not happen if images_pil was populated
            cleanup_files(temp_image_dir_path)
            return False, "Image saving failed after extraction.", None

        if len(output_image_paths) == 1:
            # Move single image out of temp dir directly into converted folder
            final_single_image_name = os.path.basename(output_image_paths[0])
            final_single_image_path = os.path.join(app.config['CONVERTED_FOLDER'], final_single_image_name)
            shutil.move(output_image_paths[0], final_single_image_path)
            cleanup_files(temp_image_dir_path) # Remove the now empty temp dir
            return True, f"PDF page converted to image: {final_single_image_name}", final_single_image_path
        else:
            zip_filename = f"{os.path.basename(output_folder_stem_for_naming)}_pages.zip"
            zip_output_path = os.path.join(app.config['CONVERTED_FOLDER'], zip_filename)
            with zipfile.ZipFile(zip_output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for img_path in output_image_paths:
                    zipf.write(img_path, os.path.basename(img_path))
            
            cleanup_files(temp_image_dir_path) # remove the temp dir with individual images
            return True, f"PDF pages converted to images and zipped into: {zip_filename}", zip_output_path

    except Exception as e:
        app.logger.error(f"PDF to Images conversion error: {e}", exc_info=True)
        cleanup_files(temp_image_dir_path if 'temp_image_dir_path' in locals() else None, *output_image_paths)
        return False, f"Failed to convert PDF to images. Ensure Poppler is installed and in PATH. Error: {str(e)}", None

def convert_pptx_to_pdf_basic(input_path, output_path): # Same
    try:
        prs = Presentation(input_path)
        # Calculate page size in points (1 inch = 72 points, EMU / 914400 = inches)
        page_width_pt = prs.slide_width.emu / 12700  # 914400 / 72 = 12700
        page_height_pt = prs.slide_height.emu / 12700
        
        c = reportlab_canvas.Canvas(output_path, pagesize=(page_width_pt, page_height_pt))

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                left_pt = shape.left.emu / 12700
                top_pt = shape.top.emu / 12700
                width_pt = shape.width.emu / 12700
                height_pt = shape.height.emu / 12700
                y_pos_reportlab = page_height_pt - top_pt - height_pt # ReportLab Y is from bottom

                if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                    text_object = c.beginText(left_pt + 5, y_pos_reportlab + height_pt - 15) # Crude baseline adjust
                    # Basic font styling attempt
                    try:
                        first_run_font = shape.text_frame.paragraphs[0].runs[0].font
                        font_name = first_run_font.name if first_run_font.name else "Helvetica"
                        font_size = first_run_font.size.pt if first_run_font.size else 10
                        text_object.setFont(font_name, font_size)
                        if first_run_font.bold: text_object.setFont(font_name+"-Bold", font_size) # Assumes bold variant exists
                        if first_run_font.italic: text_object.setFont(font_name+"-Oblique", font_size) # Assumes italic variant
                    except:
                        text_object.setFont("Helvetica", 10) # Fallback

                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_object.textLine(run.text)
                    c.drawText(text_object)

                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        temp_img_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_img_{uuid.uuid4().hex}.{image.ext}")
                        with open(temp_img_path, 'wb') as tmp_img_file:
                            tmp_img_file.write(img_bytes)
                        
                        rl_image = ReportlabImageReader(temp_img_path)
                        c.drawImage(rl_image, left_pt, y_pos_reportlab, width=width_pt, height=height_pt, preserveAspectRatio=True, mask='auto')
                        cleanup_files(temp_img_path)
                    except Exception as img_ex:
                        app.logger.warning(f"Could not process image in PPTX slide {slide_idx+1}: {img_ex}")
            c.showPage()
        c.save()
        return True, "PPTX converted to PDF with basic layout."
    except Exception as e:
        app.logger.error(f"PPTX to PDF basic conversion error: {e}", exc_info=True)
        return False, f"Failed to convert PPTX to PDF. Error: {str(e)}"

# --- Merging Functions (Same as before, ensure they are robust) ---
def merge_pdf_files_robust(input_file_paths, output_path): # Same
    merger = PdfMerger()
    try:
        for pdf_path in input_file_paths:
            merger.append(pdf_path)
        merger.write(output_path)
        return True, "PDF files merged successfully."
    finally: # Ensure merger is closed even on error
        merger.close()


def merge_docx_files_robust(input_file_paths, output_path): # Same basic append
    try:
        if not input_file_paths: return False, "No DOCX files provided."
        
        # Start with the first document to preserve its styles and sections as base
        # This is still not perfect for complex section differences.
        output_doc = DocxDocument(input_file_paths[0])
        # Remove the last empty paragraph if it exists, to avoid extra space before appending
        if output_doc.paragraphs and not output_doc.paragraphs[-1].text.strip():
             if len(output_doc.paragraphs[-1]._p.xpath(".//*")) == 0 : # truly empty
                p_to_delete = output_doc.paragraphs[-1]._p
                p_to_delete.getparent().remove(p_to_delete)


        for i in range(1, len(input_file_paths)):
            output_doc.add_page_break()
            sub_doc = DocxDocument(input_file_paths[i])
            for element in sub_doc.element.body:
                output_doc.element.body.append(element)
        
        output_doc.save(output_path)
        return True, "DOCX files merged by appending content."
    except Exception as e:
        app.logger.error(f"DOCX merge error: {e}", exc_info=True)
        return False, f"Failed to merge DOCX files. Error: {str(e)}"


# --- Flask Routes ---
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST']) # Route name for url_for: convert_file_route
def convert_file_route():
    # ... (Input validation and file saving logic from previous enhanced app.py)
    # This part remains mostly the same, ensuring it returns JSON
    if 'file_to_convert' not in request.files:
        return jsonify(success=False, error="No file part in the request."), 400
    
    file = request.files['file_to_convert']
    target_format = request.form.get('target_format')

    if file.filename == '':
        return jsonify(success=False, error="No file selected."), 400
    if not target_format:
        return jsonify(success=False, error="No target format specified."), 400

    original_filename = secure_filename(file.filename)
    input_ext = get_file_ext(original_filename)

    if input_ext not in ALLOWED_UPLOAD_EXTENSIONS:
        return jsonify(success=False, error=f"Unsupported input file type: .{input_ext}"), 400

    unique_id = uuid.uuid4().hex
    input_filename_on_server = f"{unique_id}_input.{input_ext}"
    input_filepath = os.path.join(app.config['UPLOAD_FOLDER'], input_filename_on_server)
    
    output_filename_stem = f"{unique_id}_output"
    output_ext = target_format 
    if input_ext == 'pdf' and target_format == 'images': output_ext = 'zip' # Default for multiple images
    
    output_filename_on_server = f"{output_filename_stem}.{output_ext}"
    # Default output path, might be overridden by pdf_to_images_robust
    output_filepath_final_on_server = os.path.join(app.config['CONVERTED_FOLDER'], output_filename_on_server)
    
    file.save(input_filepath)
    
    success = False
    message = "Conversion process initiated."
    processed_filename_for_download = None

    try:
        if input_ext == 'pdf':
            if target_format == 'docx':
                success, message = convert_pdf_to_docx_robust(input_filepath, output_filepath_final_on_server)
            elif target_format == 'images':
                # output_folder_stem_for_naming is used by pdf_to_images_robust to name individual images and the final zip.
                # It should be a path stem *inside* the CONVERTED_FOLDER.
                pdf_to_images_output_stem_for_naming = os.path.join(app.config['CONVERTED_FOLDER'], output_filename_stem)
                success, message, returned_path = convert_pdf_to_images_robust(input_filepath, pdf_to_images_output_stem_for_naming)
                if success and returned_path:
                    output_filepath_final_on_server = returned_path 
        elif input_ext == 'docx':
            if target_format == 'pdf':
                success, message = convert_docx_to_pdf_robust(input_filepath, output_filepath_final_on_server)
        elif input_ext == 'pptx':
            if target_format == 'pdf':
                success, message = convert_pptx_to_pdf_basic(input_filepath, output_filepath_final_on_server)
        elif input_ext in IMAGE_EXTENSIONS: # jpg, png
            if target_format == 'pdf':
                success, message = convert_images_to_pdf_robust([input_filepath], output_filepath_final_on_server)
            elif target_format in IMAGE_EXTENSIONS: # Image to Image (e.g. PNG to JPG)
                try:
                    img = Image.open(input_filepath)
                    # Handle transparency for JPG output
                    if target_format.lower() == 'jpg' and img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')
                    img.save(output_filepath_final_on_server, format=target_format.upper())
                    img.close()
                    success = True
                    message = f"Image converted to .{target_format.upper()} successfully."
                except Exception as img_e:
                    success = False
                    message = f"Image to {target_format.upper()} conversion failed: {str(img_e)}"
                    app.logger.error(f"Image to Image conversion error: {img_e}", exc_info=True)
            else:
                success = False
                message = f"Unsupported target format '{target_format}' for image input."
        else:
            success = False
            message = f"Conversion from .{input_ext} to .{target_format} is not supported."

        if success:
            processed_filename_for_download = os.path.basename(output_filepath_final_on_server)
            download_url = url_for('download_file_route', filename=processed_filename_for_download) # Corrected endpoint name
            return jsonify(success=True, message=message, filename=processed_filename_for_download, download_url=download_url)
        else:
            cleanup_files(output_filepath_final_on_server)
            return jsonify(success=False, error=message), 422

    except Exception as e:
        app.logger.error(f"General conversion route error: {e}", exc_info=True)
        cleanup_files(output_filepath_final_on_server)
        return jsonify(success=False, error=f"An unexpected server error occurred: {str(e)}"), 500
    finally:
        cleanup_files(input_filepath)


@app.route('/merge', methods=['POST']) # Route name for url_for: merge_files_route
def merge_files_route():
    # ... (Input validation and file saving logic from previous enhanced app.py)
    # This part remains mostly the same, ensuring it returns JSON
    if 'files_to_merge' not in request.files:
        return jsonify(success=False, error="No files part in the request."), 400

    files = request.files.getlist('files_to_merge')
    merge_type = request.form.get('merge_type', '').lower()

    if not files or all(f.filename == '' for f in files):
        return jsonify(success=False, error="No files selected for merging."), 400
    if not merge_type:
         return jsonify(success=False, error="Merge type not specified or detected."), 400
    if merge_type not in ['pdf', 'docx']:
        return jsonify(success=False, error=f"Merging for '.{merge_type}' files is not supported at this time."), 400 # Only PDF and DOCX for now
    if len(files) < 2:
        return jsonify(success=False, error="Please select at least two files to merge."), 400

    input_filepaths_on_server = []
    original_filenames_stem = [] # For naming the merged file
    all_files_valid_for_merge = True
    error_message_for_client = "An unknown error occurred validating files for merge."


    for file in files:
        if file and file.filename != '':
            original_filename = secure_filename(file.filename)
            file_ext = get_file_ext(original_filename)
            if file_ext != merge_type:
                all_files_valid_for_merge = False
                error_message_for_client = f"All files must be of type '.{merge_type}'. File '{original_filename}' is a '.{file_ext}'."
                break # Stop processing further files
            
            temp_input_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}_{original_filename}")
            file.save(temp_input_path)
            input_filepaths_on_server.append(temp_input_path)
            original_filenames_stem.append(os.path.splitext(original_filename)[0])
        else:
            all_files_valid_for_merge = False
            error_message_for_client = "An empty or invalid file was encountered during merge."
            break 
            
    if not all_files_valid_for_merge:
        cleanup_files(*input_filepaths_on_server) # Clean up any files already saved
        return jsonify(success=False, error=error_message_for_client), 400

    # Proceed with merging
    # Use a more descriptive name for the merged file
    merged_file_base_name = f"merged_{merge_type}_{original_filenames_stem[0]}_{uuid.uuid4().hex[:4]}"
    output_filename_on_server = f"{merged_file_base_name}.{merge_type}"
    output_filepath_on_server = os.path.join(app.config['CONVERTED_FOLDER'], output_filename_on_server)
    
    success = False
    message = "Merge process initiated." # Default message

    try:
        if merge_type == 'pdf':
            success, message = merge_pdf_files_robust(input_filepaths_on_server, output_filepath_on_server)
        elif merge_type == 'docx':
            success, message = merge_docx_files_robust(input_filepaths_on_server, output_filepath_on_server)
        # Add other merge types if implemented

        if success:
            download_url = url_for('download_file_route', filename=output_filename_on_server) # Corrected endpoint name
            return jsonify(success=True, message=message, filename=output_filename_on_server, download_url=download_url)
        else:
            cleanup_files(output_filepath_on_server) # Clean failed merge output
            return jsonify(success=False, error=message), 500 # Internal server error if merge function failed

    except Exception as e:
        app.logger.error(f"General merge route error: {e}", exc_info=True)
        cleanup_files(output_filepath_on_server) # Clean up if error
        return jsonify(success=False, error=f"An unexpected server error occurred during merge: {str(e)}"), 500
    finally:
        cleanup_files(*input_filepaths_on_server) # Always cleanup uploaded inputs


@app.route('/download/<path:filename>') # Route name for url_for: download_file_route
def download_file_route(filename): # Changed function name to match
    safe_filename = secure_filename(filename)
    if ".." in safe_filename or safe_filename.startswith("/"):
        app.logger.warning(f"Attempted directory traversal: {filename}")
        return "Invalid filename provided.", 400
    
    try:
        return send_from_directory(app.config['CONVERTED_FOLDER'], safe_filename, as_attachment=True)
    except FileNotFoundError:
        app.logger.error(f"Download error: File not found - {safe_filename}")
        return "File not found. It may have been cleaned up or never created.", 404
    # finally block for post-download cleanup was removed for simplicity;
    # files in 'converted' folder should be managed by a separate cleanup policy/task.


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)